from __future__ import annotations

import logging
from pathlib import Path
from typing import Any, Literal, cast

from langchain_community.document_loaders import PyPDFLoader

ExcelEngine = Literal["openpyxl", "xlrd"]

from .base import DocumentLoader, LoaderDependencyError, LoaderRegistry, build_document
from ..models import Document

class _ExtensionBasedLoader(DocumentLoader):
    extensions: set[str] = set()

    def supports(self, source: str | Path) -> bool:
        source_path = Path(source)
        return source_path.suffix.lower() in self.extensions


class TextFileLoader(_ExtensionBasedLoader):
    extensions = {".txt", ".md", ".rst"}

    def __init__(self, encoding: str = "utf-8") -> None:
        self._encoding = encoding

    def load(self, source: str | Path) -> list[Document]:
        path = Path(source)
        text = path.read_text(encoding=self._encoding)
        return [build_document(path, "text", text, {"extension": path.suffix.lower()})]


class PdfLoader(_ExtensionBasedLoader):
    extensions = {".pdf"}

    def load(self, source: str | Path) -> list[Document]:
        path = Path(source)
        try:
            docs = PyPDFLoader(str(path)).load()
        except Exception as exc:
            logging.warning(f"Failed to load PDF document {path}: {exc}")
            return []

        return [
            build_document(
                path,
                "pdf",
                doc.page_content,
                {
                    "extension": ".pdf",
                    "page": doc.metadata.get("page"),
                    "source": doc.metadata.get("source", str(path)),
                },
            )
            for doc in docs
            if doc.page_content
        ]


class WordDocxLoader(_ExtensionBasedLoader):
    extensions = {".docx"}

    def load(self, source: str | Path) -> list[Document]:
        try:
            import docx2txt
        except ImportError as exc:
            raise LoaderDependencyError(
                "Word support requires `docx2txt`. Install it with: uv add docx2txt"
            ) from exc

        path = Path(source)
        try:
            text = docx2txt.process(str(path))
            return [build_document(path, "word", text, {"extension": ".docx"})]
        except Exception as exc:
            logging.warning(f"Failed to load Word document {path}: {exc}")
            return []
        
class WordDocLoader(_ExtensionBasedLoader):
    """Loader for legacy binary ``.doc`` files (pre-2007 Word format).

    Converts ``.doc`` to ``.docx`` via a headless LibreOffice (``soffice``)
    invocation in an isolated user profile (to avoid profile-lock collisions
    under concurrent Celery workers / a running LibreOffice GUI), then
    extracts text with ``docx2txt``.
    """

    extensions = {".doc"}

    def load(self, source: str | Path) -> list[Document]:
        import shutil
        import subprocess
        import tempfile

        soffice = shutil.which("soffice") or shutil.which("libreoffice")
        if not soffice:
            raise LoaderDependencyError(
                "Legacy .doc support requires LibreOffice (`soffice`) on PATH. "
                "Install with: brew install --cask libreoffice"
            )

        try:
            import docx2txt
        except ImportError as exc:
            raise LoaderDependencyError(
                "Legacy .doc support requires `docx2txt`. Install it with: uv add docx2txt"
            ) from exc

        path = Path(source)
        with tempfile.TemporaryDirectory() as tmpdir:
            tmp = Path(tmpdir)
            profile_uri = (tmp / "profile").as_uri()
            outdir = tmp / "out"
            outdir.mkdir()
            cmd = [
                soffice,
                f"-env:UserInstallation={profile_uri}",
                "--headless",
                "--norestore",
                "--nologo",
                "--nofirststartwizard",
                "--convert-to",
                "docx:MS Word 2007 XML",
                "--outdir",
                str(outdir),
                str(path),
            ]
            try:
                result = subprocess.run(
                    cmd,
                    capture_output=True,
                    text=True,
                    timeout=120,
                    check=False,
                )
            except (subprocess.TimeoutExpired, OSError) as exc:
                logging.warning(f"soffice invocation failed for {path}: {exc}")
                return []

            converted = outdir / (path.stem + ".docx")
            if result.returncode != 0 or not converted.exists():
                logging.warning(
                    "soffice failed to convert %s to docx (exit=%s). "
                    "stdout=%r stderr=%r",
                    path, result.returncode, result.stdout, result.stderr,
                )
                return []

            try:
                text = docx2txt.process(str(converted))
            except Exception as exc:
                logging.warning(f"Failed to extract text from converted {converted}: {exc}")
                return []

        if not text or not text.strip():
            return []
        return [build_document(path, "word", text, {"extension": ".doc"})]


class ExcelLoader(_ExtensionBasedLoader):
    extensions = {".xlsx", ".xls"}

    def load(self, source: str | Path) -> list[Document]:
        try:
            import pandas as pd
        except ImportError as exc:
            raise LoaderDependencyError(
                "Excel support requires `pandas`. Install it with: uv add pandas"
            ) from exc

        path = Path(source)
        engine = self._engine_for_suffix(path.suffix.lower())
        if engine == "xlrd":
            try:
                import xlrd  # noqa: F401
            except ImportError as exc:
                raise LoaderDependencyError(
                    "Legacy .xls support requires `xlrd`. Install it with: uv add xlrd"
                ) from exc

        try:
            chunks: list[str] = []
            workbook: dict[str, "pd.DataFrame"] = pd.read_excel(
                path,
                sheet_name=None,
                header=None,
                engine=engine,
                dtype=object,  # type: ignore[arg-type]
            )

            for sheet_name, frame in workbook.items():
                rows: list[str] = []
                for row in frame.itertuples(index=False, name=None):
                    line = "\t".join(self._stringify_cell(cell) for cell in row).strip()
                    if line:
                        rows.append(line)
                if rows:
                    chunks.append(f"# Sheet: {sheet_name}\n" + "\n".join(rows))

            text = "\n\n".join(chunks)
            return [build_document(path, "excel", text, {"extension": path.suffix.lower()})]
        except Exception as exc:
            logging.warning(f"Failed to load Excel document {path}: {exc}")
            return []

    @staticmethod
    def _engine_for_suffix(suffix: str) -> ExcelEngine:
        if suffix == ".xls":
            return "xlrd"
        return "openpyxl"

    @staticmethod
    def _stringify_cell(value: object) -> str:
        try:
            import pandas as pd
        except ImportError:
            return "" if value is None else str(value)

        if value is None or bool(pd.isna(cast(Any, value))):
            return ""
        return str(value)


class FileSystemLoader:
    """Loads documents from a file or directory, dispatching by file extension."""

    def __init__(self, path: str | Path, recursive: bool = False, registry: LoaderRegistry | None = None) -> None:
        self._path = Path(path)
        self._recursive = recursive
        self._registry = registry

    def load(self) -> list[Document]:
        if self._path.is_file():
            return self._load_file(self._path)
        if self._path.is_dir():
            return self._load_directory(self._path)
        raise ValueError(f"Path does not exist or is not a file/directory: {self._path}")

    def _load_file(self, path: Path) -> list[Document]:
        from . import build_default_registry
        registry = self._registry or build_default_registry()
        try:
            return registry.load(path)
        except ValueError:
            return []

    def _load_directory(self, directory: Path) -> list[Document]:
        pattern = "**/*" if self._recursive else "*"
        documents: list[Document] = []
        for entry in sorted(directory.glob(pattern)):
            if entry.is_file():
                documents.extend(self._load_file(entry))
        return documents


class PowerPointLoader(_ExtensionBasedLoader):
    extensions = {".pptx"}

    def load(self, source: str | Path) -> list[Document]:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise LoaderDependencyError(
                "PowerPoint support requires `python-pptx`. Install it with: uv add python-pptx"
            ) from exc

        path = Path(source)
        presentation = Presentation(str(path))

        slide_chunks: list[str] = []
        for idx, slide in enumerate(presentation.slides, start=1):
            lines: list[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text:
                    stripped = text.strip()
                    if stripped:
                        lines.append(stripped)
            if lines:
                slide_chunks.append(f"# Slide {idx}\n" + "\n".join(lines))

        text = "\n\n".join(slide_chunks)
        return [build_document(path, "powerpoint", text, {"extension": ".pptx"})]

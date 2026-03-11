from __future__ import annotations

from pathlib import Path

from langchain_community.document_loaders import PyPDFLoader

from .base import DocumentLoader, LoaderDependencyError, LoaderRegistry, build_document
from ..models import Document


class _ExtensionBasedLoader(DocumentLoader):
    extensions: set[str] = set()

    def supports(self, source: str | Path) -> bool:
        source_path = Path(source)
        return source_path.suffix.lower() in self.extensions


class TextFileLoader(_ExtensionBasedLoader):
    extensions = {".txt", ".md", ".rst"}

    def __init__(self, encoding: str = "utf-8") -> None:
        self._encoding = encoding

    def load(self, source: str | Path) -> list[Document]:
        path = Path(source)
        text = path.read_text(encoding=self._encoding)
        return [build_document(path, "text", text, {"extension": path.suffix.lower()})]


class PdfLoader(_ExtensionBasedLoader):
    extensions = {".pdf"}

    def load(self, source: str | Path) -> list[Document]:
        path = Path(source)
        docs = PyPDFLoader(str(path)).load()
        return [
            build_document(
                path,
                "pdf",
                doc.page_content,
                {
                    "extension": ".pdf",
                    "page": doc.metadata.get("page"),
                    "source": doc.metadata.get("source", str(path)),
                },
            )
            for doc in docs
            if doc.page_content
        ]


class WordLoader(_ExtensionBasedLoader):
    extensions = {".docx"}

    def load(self, source: str | Path) -> list[Document]:
        try:
            import docx2txt
        except ImportError as exc:
            raise LoaderDependencyError(
                "Word support requires `docx2txt`. Install it with: uv add docx2txt"
            ) from exc

        path = Path(source)
        text = docx2txt.process(str(path))
        return [build_document(path, "word", text, {"extension": ".docx"})]


class ExcelLoader(_ExtensionBasedLoader):
    extensions = {".xlsx", ".xls"}

    def load(self, source: str | Path) -> list[Document]:
        try:
            from openpyxl import load_workbook
        except ImportError as exc:
            raise LoaderDependencyError(
                "Excel support requires `openpyxl`. Install it with: uv add openpyxl"
            ) from exc

        path = Path(source)
        wb = load_workbook(path, read_only=True, data_only=True)

        chunks: list[str] = []
        for sheet in wb.worksheets:
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                line = "\t".join("" if cell is None else str(cell) for cell in row).strip()
                if line:
                    rows.append(line)
            if rows:
                chunks.append(f"# Sheet: {sheet.title}\n" + "\n".join(rows))

        text = "\n\n".join(chunks)
        return [build_document(path, "excel", text, {"extension": path.suffix.lower()})]


class FileSystemLoader:
    """Loads documents from a file or directory, dispatching by file extension."""

    def __init__(self, path: str | Path, recursive: bool = False, registry: LoaderRegistry | None = None) -> None:
        self._path = Path(path)
        self._recursive = recursive
        self._registry = registry

    def load(self) -> list[Document]:
        if self._path.is_file():
            return self._load_file(self._path)
        if self._path.is_dir():
            return self._load_directory(self._path)
        raise ValueError(f"Path does not exist or is not a file/directory: {self._path}")

    def _load_file(self, path: Path) -> list[Document]:
        from . import build_default_registry
        registry = self._registry or build_default_registry()
        try:
            return registry.load(path)
        except ValueError:
            return []

    def _load_directory(self, directory: Path) -> list[Document]:
        pattern = "**/*" if self._recursive else "*"
        documents: list[Document] = []
        for entry in sorted(directory.glob(pattern)):
            if entry.is_file():
                documents.extend(self._load_file(entry))
        return documents


class PowerPointLoader(_ExtensionBasedLoader):
    extensions = {".pptx"}

    def load(self, source: str | Path) -> list[Document]:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise LoaderDependencyError(
                "PowerPoint support requires `python-pptx`. Install it with: uv add python-pptx"
            ) from exc

        path = Path(source)
        presentation = Presentation(str(path))

        slide_chunks: list[str] = []
        for idx, slide in enumerate(presentation.slides, start=1):
            lines: list[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text:
                    stripped = text.strip()
                    if stripped:
                        lines.append(stripped)
            if lines:
                slide_chunks.append(f"# Slide {idx}\n" + "\n".join(lines))

        text = "\n\n".join(slide_chunks)
        return [build_document(path, "powerpoint", text, {"extension": ".pptx"})]
